#!/usr/bin/env python3
"""Multi-format source ingester for the portfolio content pipeline.

Turns raw project source material (PDF / PPTX / images / text) into a clean,
agent-friendly *staging* folder:

    content/_staging/<slug>/
        text.md          concatenated text, with per-source markers
        images/          deduped, EXIF-stripped, width-capped images
        manifest.json    structured index of every extracted asset

An authoring agent then reads ``text.md`` + ``manifest.json`` and writes
``src/content/projects/<slug>.mdx`` (see docs/CONTENT_PIPELINE.md).

Usage:
    python scripts/ingest.py <slug>
    python scripts/ingest.py <slug> --src some/folder --out some/out --max-width 2000

Dependencies (see scripts/requirements.txt) are only required for the formats
you actually use; the script reports a helpful message if one is missing.
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import sys
from dataclasses import dataclass, field
from pathlib import Path

# --- Configuration ----------------------------------------------------------

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}
TEXT_EXTS = {".txt", ".md", ".markdown"}
PDF_EXTS = {".pdf"}
PPTX_EXTS = {".pptx", ".ppt"}
DOCX_EXTS = {".docx"}

DEFAULT_MAX_WIDTH = 2400
DEFAULT_JPEG_QUALITY = 82
# Ignore tiny images (logos, bullet icons, line decorations) that add no value.
MIN_DIMENSION = 80


# --- Data model -------------------------------------------------------------


@dataclass
class Asset:
    file: str  # path relative to the staging images/ dir
    source: str  # originating source filename
    location: str  # "page 3", "slide 2", or "file"
    width: int
    height: int
    bytes: int
    sha256: str


@dataclass
class IngestResult:
    text_blocks: list[str] = field(default_factory=list)
    assets: list[Asset] = field(default_factory=list)
    skipped: list[str] = field(default_factory=list)


# --- Helpers ----------------------------------------------------------------


def _missing_dep(pkg: str, fmt: str) -> SystemExit:
    return SystemExit(
        f"Cannot process {fmt} files: the '{pkg}' package is not installed.\n"
        f"Install the pipeline dependencies with:\n"
        f"    pip install -r scripts/requirements.txt"
    )


def slugify_stem(name: str) -> str:
    keep = [c if c.isalnum() else "-" for c in name.lower()]
    out = "".join(keep)
    while "--" in out:
        out = out.replace("--", "-")
    return out.strip("-") or "img"


class ImageWriter:
    """Writes deduped, normalised images into the staging images/ folder."""

    def __init__(self, out_dir: Path, max_width: int, quality: int) -> None:
        self.out_dir = out_dir
        self.max_width = max_width
        self.quality = quality
        self._seen_hashes: dict[str, str] = {}  # sha256 -> written filename
        self._used_names: set[str] = set()
        try:
            from PIL import Image  # noqa: F401
        except ImportError as exc:  # pragma: no cover
            raise _missing_dep("Pillow", "image") from exc

    def _unique_name(self, stem: str, ext: str) -> str:
        candidate = f"{stem}{ext}"
        i = 1
        while candidate in self._used_names:
            candidate = f"{stem}-{i}{ext}"
            i += 1
        self._used_names.add(candidate)
        return candidate

    def write(self, blob: bytes, stem: str, source: str, location: str) -> Asset | None:
        from PIL import Image

        sha = hashlib.sha256(blob).hexdigest()
        if sha in self._seen_hashes:
            return None  # exact duplicate already written

        try:
            img = Image.open(io.BytesIO(blob))
            img.load()
        except Exception:
            return None  # not a decodable image

        if img.width < MIN_DIMENSION and img.height < MIN_DIMENSION:
            return None

        has_alpha = img.mode in ("RGBA", "LA", "P") and "transparency" in img.info
        fmt = "PNG" if (has_alpha or img.mode == "P") else "JPEG"
        ext = ".png" if fmt == "PNG" else ".jpg"

        if fmt == "JPEG" and img.mode != "RGB":
            img = img.convert("RGB")
        elif fmt == "PNG" and img.mode == "P":
            img = img.convert("RGBA")

        if img.width > self.max_width:
            ratio = self.max_width / img.width
            img = img.resize((self.max_width, round(img.height * ratio)), Image.LANCZOS)

        name = self._unique_name(slugify_stem(stem), ext)
        dest = self.out_dir / name
        save_kwargs = {"optimize": True}
        if fmt == "JPEG":
            save_kwargs.update(quality=self.quality, progressive=True)
        img.save(dest, fmt, **save_kwargs)  # re-encode strips EXIF/metadata

        self._seen_hashes[sha] = name
        return Asset(
            file=name,
            source=source,
            location=location,
            width=img.width,
            height=img.height,
            bytes=dest.stat().st_size,
            sha256=sha,
        )


# --- Per-format extractors --------------------------------------------------


def ingest_pdf(path: Path, writer: ImageWriter, result: IngestResult) -> None:
    try:
        import fitz  # PyMuPDF
    except ImportError as exc:
        raise _missing_dep("pymupdf", "PDF") from exc

    doc = fitz.open(path)
    stem = path.stem
    for page_index in range(len(doc)):
        page = doc[page_index]
        page_no = page_index + 1
        text = page.get_text("text").strip()
        if text:
            result.text_blocks.append(f"### {path.name} — page {page_no}\n\n{text}")
        for img_index, info in enumerate(page.get_images(full=True), start=1):
            xref = info[0]
            try:
                extracted = doc.extract_image(xref)
            except Exception:
                continue
            asset = writer.write(
                extracted["image"],
                stem=f"{stem}-p{page_no}-{img_index}",
                source=path.name,
                location=f"page {page_no}",
            )
            if asset:
                result.assets.append(asset)
    doc.close()


def ingest_pptx(path: Path, writer: ImageWriter, result: IngestResult) -> None:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise _missing_dep("python-pptx", "PowerPoint") from exc

    prs = Presentation(str(path))
    stem = path.stem
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        img_index = 0

        def walk(shapes) -> None:
            nonlocal img_index
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shape.shapes)
                    continue
                if shape.has_text_frame and shape.text_frame.text.strip():
                    texts.append(shape.text_frame.text.strip())
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_index += 1
                    try:
                        image = shape.image
                    except Exception:
                        continue
                    asset = writer.write(
                        image.blob,
                        stem=f"{stem}-s{slide_index}-{img_index}",
                        source=path.name,
                        location=f"slide {slide_index}",
                    )
                    if asset:
                        result.assets.append(asset)

        walk(slide.shapes)
        if texts:
            joined = "\n\n".join(texts)
            result.text_blocks.append(f"### {path.name} — slide {slide_index}\n\n{joined}")


def ingest_image(path: Path, writer: ImageWriter, result: IngestResult) -> None:
    asset = writer.write(
        path.read_bytes(),
        stem=path.stem,
        source=path.name,
        location="file",
    )
    if asset:
        result.assets.append(asset)
    else:
        result.skipped.append(f"{path.name} (duplicate or too small)")


def ingest_text(path: Path, _writer: ImageWriter, result: IngestResult) -> None:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    if text:
        result.text_blocks.append(f"### {path.name}\n\n{text}")


def ingest_docx(path: Path, writer: ImageWriter, result: IngestResult) -> None:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise _missing_dep("python-docx", "Word") from exc

    document = docx.Document(str(path))
    paras = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    if paras:
        result.text_blocks.append(f"### {path.name}\n\n" + "\n\n".join(paras))
    stem = path.stem
    for i, rel in enumerate(document.part.rels.values(), start=1):
        if "image" not in rel.reltype:
            continue
        try:
            blob = rel.target_part.blob
        except Exception:
            continue
        asset = writer.write(blob, stem=f"{stem}-{i}", source=path.name, location="file")
        if asset:
            result.assets.append(asset)


DISPATCH = [
    (PDF_EXTS, ingest_pdf),
    (PPTX_EXTS, ingest_pptx),
    (DOCX_EXTS, ingest_docx),
    (IMAGE_EXTS, ingest_image),
    (TEXT_EXTS, ingest_text),
]


# --- Orchestration ----------------------------------------------------------


def collect_sources(src_dir: Path) -> list[Path]:
    files = [p for p in sorted(src_dir.rglob("*")) if p.is_file() and not p.name.startswith(".")]
    return files


def run(slug: str, src: Path, out: Path, max_width: int, quality: int) -> int:
    if not src.exists():
        print(f"✗ Source folder not found: {src}", file=sys.stderr)
        print(f"  Create it and drop your source files in, then re-run.", file=sys.stderr)
        return 1

    sources = collect_sources(src)
    if not sources:
        print(f"✗ No source files found in {src}", file=sys.stderr)
        return 1

    images_dir = out / "images"
    images_dir.mkdir(parents=True, exist_ok=True)
    writer = ImageWriter(images_dir, max_width=max_width, quality=quality)
    result = IngestResult()

    print(f"Ingesting '{slug}' from {src}")
    for path in sources:
        ext = path.suffix.lower()
        handler = next((fn for exts, fn in DISPATCH if ext in exts), None)
        if handler is None:
            result.skipped.append(f"{path.name} (unsupported type)")
            print(f"  · skip {path.name} (unsupported)")
            continue
        print(f"  → {path.name}")
        handler(path, writer, result)

    # Write text.md
    header = (
        f"# Source material — {slug}\n\n"
        f"> Auto-generated by scripts/ingest.py. Read this plus images/manifest.json,\n"
        f"> then author src/content/projects/{slug}.mdx per docs/CONTENT_PIPELINE.md.\n"
    )
    body = "\n\n---\n\n".join(result.text_blocks) if result.text_blocks else "_No text extracted._"
    (out / "text.md").write_text(f"{header}\n{body}\n", encoding="utf-8")

    # Write manifest.json
    manifest = {
        "slug": slug,
        "source_dir": str(src),
        "max_width": max_width,
        "counts": {
            "sources": len(sources),
            "images": len(result.assets),
            "text_blocks": len(result.text_blocks),
            "skipped": len(result.skipped),
        },
        "images": [asset.__dict__ for asset in result.assets],
        "skipped": result.skipped,
    }
    (out / "manifest.json").write_text(json.dumps(manifest, indent=2), encoding="utf-8")

    print(
        f"\n✓ Done. {len(result.assets)} image(s), {len(result.text_blocks)} text block(s).\n"
        f"  Staging: {out}\n"
        f"  Next: author src/content/projects/{slug}.mdx (see docs/CONTENT_PIPELINE.md)."
    )
    if result.skipped:
        print(f"  Skipped {len(result.skipped)}: " + ", ".join(result.skipped[:8]))
    return 0


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Ingest raw source files into the content staging folder.",
    )
    parser.add_argument("slug", help="Project slug, e.g. 'autonomous-drone-mav'")
    parser.add_argument("--src", type=Path, help="Source dir (default content/_sources/<slug>)")
    parser.add_argument("--out", type=Path, help="Output dir (default content/_staging/<slug>)")
    parser.add_argument("--max-width", type=int, default=DEFAULT_MAX_WIDTH, help="Cap image width")
    parser.add_argument("--quality", type=int, default=DEFAULT_JPEG_QUALITY, help="JPEG quality")
    args = parser.parse_args(argv)

    repo_root = Path(__file__).resolve().parent.parent
    src = args.src or repo_root / "content" / "_sources" / args.slug
    out = args.out or repo_root / "content" / "_staging" / args.slug
    return run(args.slug, src, out, args.max_width, args.quality)


if __name__ == "__main__":
    raise SystemExit(main())
